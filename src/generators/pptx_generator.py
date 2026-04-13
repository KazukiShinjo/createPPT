"""PowerPoint生成エンジン"""

import logging
from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from src.models.presentation import Presentation
from src.models.styles import Styles


logger = logging.getLogger(__name__)


class PresentationGenerator:
    """PowerPointスライドを生成する"""

    def __init__(self):
        self.styles = Styles()
    
    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """16進カラーコードをRGBColorに変換"""
        hex_color = hex_color.lstrip('#')
        return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))

    def generate(self, presentation: Presentation, output_path: str) -> None:
        """
        プレゼンテーションオブジェクトからPowerPointファイルを生成
        
        Args:
            presentation: Presentationオブジェクト
            output_path: 出力ファイルパス
            
        Raises:
            IOError: ファイル書き込みに失敗した場合
        """
        try:
            prs = PptxPresentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)

            # タイトルスライドを追加
            self._add_title_slide(prs, presentation.title)

            # コンテンツスライドを追加
            for slide_data in presentation.slides:
                self._add_content_slide(prs, slide_data.title, slide_data.content)

            # ファイルに保存
            prs.save(output_path)
            logger.info(f"✅ PowerPointファイルを生成しました: {output_path}")

        except IOError as e:
            logger.error(f"ファイル書き込みエラー: {e}")
            raise IOError(f"ファイル書き込みに失敗しました。ディスク容量と権限を確認してください: {e}")
        except Exception as e:
            logger.error(f"PowerPoint生成エラー: {e}")
            raise

    def _add_title_slide(self, prs: PptxPresentation, title: str) -> None:
        """タイトルスライドを追加"""
        slide_layout = prs.slide_layouts[0]  # タイトルスライドレイアウト
        slide = prs.slides.add_slide(slide_layout)

        title_shape = slide.shapes.title
        title_shape.text = title

        # スタイル適用
        title_frame = title_shape.text_frame
        for paragraph in title_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(self.styles.sizes.title_slide)
                run.font.color.rgb = self._hex_to_rgb(self.styles.colors.heading)

    def _add_content_slide(self, prs: PptxPresentation, title: str, content: list) -> None:
        """コンテンツスライドを追加"""
        slide_layout = prs.slide_layouts[1]  # タイトル・コンテンツレイアウト
        slide = prs.slides.add_slide(slide_layout)

        # タイトルを設定
        title_shape = slide.shapes.title
        title_shape.text = title

        title_frame = title_shape.text_frame
        for paragraph in title_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(self.styles.sizes.slide_title)
                run.font.color.rgb = self._hex_to_rgb(self.styles.colors.heading)

        # コンテンツを設定
        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.clear()

        for idx, item in enumerate(content):
            if idx == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.text = item
            p.level = 0
            p.font.size = Pt(self.styles.sizes.body)
            p.font.color.rgb = self._hex_to_rgb(self.styles.colors.text)
